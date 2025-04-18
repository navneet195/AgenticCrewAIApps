def render_tab1():
    import streamlit as st
    import pandas as pd
    from PyPDF2 import PdfReader
    from pptx import Presentation
    from crewai import Agent, Task, Crew, Process, LLM
    from crewai.knowledge.source.string_knowledge_source import StringKnowledgeSource
    from config import api_key  # Ensure this is securely loaded
    from os import environ as env
    from config import llm

    from dotenv import load_dotenv
    load_dotenv()

    from openai import OpenAI

    api_key = env["OPENAI_API_KEY"]
    client = OpenAI(api_key=api_key)

    # ===== Streamlit Setup =====
    # Initialize session state
    if 'chat_history' not in st.session_state:
        st.session_state.chat_history = []

    if 'file_content' not in st.session_state:
        st.session_state.file_content = None


    # ===== Display Chat History =====
    def display_chat_history():
        for role, message in st.session_state.chat_history:
            label = "**User:**" if role == "User" else "**AI:**"
            st.markdown(f"{label} {message}")


    # ===== File Extraction Utils =====
    def extract_pdf(file):
        try:
            reader = PdfReader(file)
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception as e:
            st.error(f"PDF Extraction Error: {e}")
            return "Unable to extract PDF text."


    def extract_ppt(file):
        try:
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            st.error(f"PPT Extraction Error: {e}")
            return "Unable to extract PPT text."


    def extract_csv(file):
        try:
            df = pd.read_csv(file)
            st.write(df)
            return df.to_csv(index=False)
        except Exception as e:
            st.error(f"CSV Extraction Error: {e}")
            return "Unable to extract CSV file."


    def process_uploaded_file(file):
        file_type = file.type
        if file_type == "application/pdf":
            return extract_pdf(file)
        elif file_type in [
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ]:
            return extract_ppt(file)
        elif file_type == "text/csv":
            return extract_csv(file)
        else:
            return "Unsupported file type."


    # ===== Main App Interface =====
    st.header("📤 Upload and Extract File")

    uploaded_files = st.file_uploader(
        "Choose multiple files (PDF, PPTX, CSV)", 
        type=["pdf", "ppt", "pptx", "csv"], 
        accept_multiple_files=True
    )

    pdf_ppt_content = ""
    csv_content = ""

    for file in uploaded_files:
        if file.type == "application/pdf":
            pdf_ppt_content += extract_pdf(file)
        elif file.type in [
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ]:
            pdf_ppt_content += extract_ppt(file)
        elif file.type == "text/csv":
            csv_content += extract_csv(file)


    if uploaded_files:
        combined_content = pdf_ppt_content + csv_content
        st.session_state.file_content = combined_content
        st.subheader("📚 Combined Extracted Content")
        st.text_area("Combined Content", combined_content, height=300)

    # Guard clause if no content
    if not st.session_state.file_content:
        st.warning("Please upload a file first.")
        st.stop()


    # ===== Knowledge & LLM Setup =====
    if st.session_state.file_content:
        pdf_ppt_source = StringKnowledgeSource(content=pdf_ppt_content)
        csv_source = StringKnowledgeSource(content=csv_content)

        my_source = StringKnowledgeSource(
            content=my_text,
            embedding_model="cohere-embed-v3",
            vector_store="chroma"
        )


    def generate_sample_questions(context, num_questions=5):
        prompt = f"""
        Based on the following document content, generate {num_questions} useful and diverse questions a user might ask:

        ---
        {context[:4000]}
        ---
        Format as a numbered list.
        """

        response = client.chat.completions.create(
            model="gpt-4o-mini",  # or "gpt-3.5-turbo", gpt-4 if you prefer
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=300,
        )

        reply = response.choices[0].message.content
        # Parse and return question list
        return [line.strip("0123456789. ") for line in reply.split("\n") if line.strip()]


    pdf_ppt_agent = Agent(
        role="Document Analyst",
        goal="Answer questions from PDFs and PPTs",
        backstory="Expert in analyzing and summarizing documents",
        llm=llm,
        verbose=True,
    )

    csv_agent = Agent(
        role="Data Analyst",
        goal="Answer questions from CSV data",
        backstory="Skilled at reading and interpreting structured data tables",
        llm=llm,
        verbose=True,
    )

    task_for_doc = Task(
        description="Answer user's question using information from PDFs and PPTs",
        expected_output="Insightful summary or answer from documents",
        agent=pdf_ppt_agent
    )

    task_for_csv = Task(
        description="Answer user's question using the CSV data",
        expected_output="Answer based on data analysis or trends",
        agent=csv_agent
    )

    # Create the crew
    crew = Crew(
        agents=[pdf_ppt_agent, csv_agent],
        tasks=[task_for_doc, task_for_csv],
        knowledge_sources=[pdf_ppt_source, csv_source],
        process=Process.sequential,
        manager_llm=llm,
        verbose=True
    )

    # ===== Chat UI =====
    st.divider()
    st.subheader("💬 Ask a Question")

    # Show suggestions based on context
    if st.session_state.file_content:
        st.markdown("💡 Suggested Questions:")
        sample_questions = generate_sample_questions(st.session_state.file_content)
        for q in sample_questions:
            st.markdown(f"- {q}")

    display_chat_history()
    user_query = st.text_input("Enter your question", key="docs_chatbot_input")

    if st.button("Send") and user_query.strip():
        with st.spinner("Thinking..."):
            results = crew.kickoff(inputs={"question": user_query})

            st.session_state.chat_history.append(("User", user_query))
            st.session_state.chat_history.append(("AI", results))

            st.success("Response received!")
            display_chat_history()

            feedback = st.radio("Was the answer helpful?", ["Yes", "No"], horizontal=True)
            if feedback:
                st.write(f"Thanks for your feedback: {feedback}")

            for file in uploaded_files:
                st.download_button(f"Download {file.name}", data=file, file_name=file.name)
