import streamlit as st
import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation
from crewai import Agent, Task, Crew, Process, LLM
from crewai.knowledge.source.string_knowledge_source import StringKnowledgeSource
from config import api_key  # Ensure this is securely loaded

# ===== Streamlit Setup =====
st.set_page_config(page_title="📚 Document Handling with RAG", page_icon="📄")

# Initialize session state
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []

if 'file_content' not in st.session_state:
    st.session_state.file_content = None


# ===== Display Chat History =====
def display_chat_history():
    for role, message in st.session_state.chat_history:
        label = "**User:**" if role == "User" else "**AI:**"
        st.markdown(f"{label} {message}")


# ===== File Extraction Utils =====
def extract_pdf(file):
    try:
        reader = PdfReader(file)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as e:
        st.error(f"PDF Extraction Error: {e}")
        return "Unable to extract PDF text."


def extract_ppt(file):
    try:
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"PPT Extraction Error: {e}")
        return "Unable to extract PPT text."


def extract_csv(file):
    try:
        df = pd.read_csv(file)
        st.write(df)
        return df.to_csv(index=False)
    except Exception as e:
        st.error(f"CSV Extraction Error: {e}")
        return "Unable to extract CSV file."


def process_uploaded_file(file):
    file_type = file.type
    if file_type == "application/pdf":
        return extract_pdf(file)
    elif file_type in [
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ]:
        return extract_ppt(file)
    elif file_type == "text/csv":
        return extract_csv(file)
    else:
        return "Unsupported file type."


# ===== Main App Interface =====
st.header("📤 Upload and Extract File")

uploaded_files = st.file_uploader(
    "Choose multiple files (PDF, PPTX, CSV)", 
    type=["pdf", "ppt", "pptx", "csv"], 
    accept_multiple_files=True
)

def extract_and_combine_contents(files):
    all_text = ""
    for file in files:
        file_type = file.type
        if file_type == "application/pdf":
            content = extract_pdf(file)
        elif file_type in [
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ]:
            content = extract_ppt(file)
        elif file_type == "text/csv":
            content = extract_csv(file)
        else:
            content = "Unsupported file type."

        # Add file name as context (optional)
        all_text += f"\n---\n📄 File: {file.name}\n{content}\n"
    return all_text


if uploaded_files:
    combined_content = extract_and_combine_contents(uploaded_files)
    st.session_state.file_content = combined_content

    st.subheader("📚 Combined Extracted Content")
    st.text_area("Combined Content", combined_content, height=300)

# Guard clause if no content
if not st.session_state.file_content:
    st.warning("Please upload a file first.")
    st.stop()


# ===== Knowledge & LLM Setup =====
if st.session_state.file_content:
    source = StringKnowledgeSource(
        content=st.session_state.file_content,
        chunk_size=4000,
        chunk_overlap=200
    )

llm = LLM(
    model="openai/gpt-4o-mini",  # You can change to gpt-4, gpt-3.5, etc.
    temperature=0.8,
    max_tokens=150,
    top_p=0.9,
    frequency_penalty=0.1,
    presence_penalty=0.1,
    stop=["END"],
    seed=42
)

agent = Agent(
    role="Document Q&A Agent",
    goal="Answer questions based on provided documents.",
    backstory="You are a helpful assistant who provides context-aware answers based on uploaded documents.",
    verbose=True,
    allow_delegation=False,
    llm=llm,
)

task = Task(
    description="Answer the user's question based on the provided document.",
    expected_output="A concise and relevant answer.",
    agent=agent,
)

crew = Crew(
    agents=[agent],
    tasks=[task],
    process=Process.sequential,
    verbose=True,
    knowledge_sources=[source]
)


# ===== Chat UI =====
st.divider()
st.subheader("💬 Ask a Question")
display_chat_history()

user_query = st.text_input("Your Message")

if st.button("Send") and user_query.strip():
    with st.spinner("Thinking..."):

        # Run crew and get answer
        results = crew.kickoff(inputs={"question": user_query})
        
        # Optional: inspect results if it's structured
        st.session_state.chat_history.append(("User", user_query))
        st.session_state.chat_history.append(("AI", results))

        st.success("Response received!")
        display_chat_history()

        # Feedback
        feedback = st.radio("Was the answer helpful?", ["Yes", "No"], horizontal=True)
        if feedback:
            st.write(f"Thanks for your feedback: {feedback}")

        # Download original file
        for file in uploaded_files:
            st.download_button(f"Download {file.name}", data=file, file_name=file.name)
