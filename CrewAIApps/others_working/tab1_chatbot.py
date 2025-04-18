def render_tab1():
    import streamlit as st
    import pandas as pd
    from PyPDF2 import PdfReader
    from pptx import Presentation
    from crewai import Agent, Task, Crew, Process
    from crewai.knowledge.source.string_knowledge_source import StringKnowledgeSource
    from dotenv import load_dotenv
    from os import environ as env
    from openai import OpenAI
    from config import llm

    load_dotenv()
    client = OpenAI(api_key=env["OPENAI_API_KEY"])

    # Chunking helper
    def chunk_text(text, max_chunk_length=2000):
        words = text.split()
        chunks, chunk = [], []
        length = 0
        for word in words:
            if length + len(word) <= max_chunk_length:
                chunk.append(word)
                length += len(word) + 1
            else:
                chunks.append(" ".join(chunk))
                chunk = [word]
                length = len(word)
        if chunk:
            chunks.append(" ".join(chunk))
        return chunks

    # Session state init
    if 'chat_history' not in st.session_state:
        st.session_state.chat_history = []

    st.header("📤 Upload PDF / PPT / CSV")

    uploaded_files = st.file_uploader(
        "Upload any of the following files",
        type=["pdf", "ppt", "pptx", "csv"],
        accept_multiple_files=True
    )

    file_text = {"pdf_ppt": "", "csv": "", "csv_df": None}

    # Extract file contents
    for file in uploaded_files:
        if file.type == "application/pdf":
            reader = PdfReader(file)
            pdf_text = "\n".join(page.extract_text() or "" for page in reader.pages)
            file_text["pdf_ppt"] += pdf_text + "\n"

        elif file.type in [
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ]:
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        file_text["pdf_ppt"] += shape.text + "\n"

        elif file.type == "text/csv":
            df = pd.read_csv(file)
            file_text["csv_df"] = df
            file_text["csv"] += df.to_csv(index=False)

    # Show content preview
    full_text = file_text["pdf_ppt"] + file_text["csv"]
    if not full_text.strip():
        st.warning("Please upload a supported file.")
        st.stop()

    st.subheader("📑 Preview of Extracted Content")
    st.text_area("Preview", full_text[:3000], height=300)

    # === Build agent-specific knowledge ===
    knowledge_sources = []

    if file_text["pdf_ppt"]:
        pdf_chunks = chunk_text(file_text["pdf_ppt"])
        pdf_text_combined = "\n---\n".join(pdf_chunks[:10])  # limit chunk count
        pdf_source = StringKnowledgeSource(content=pdf_text_combined)
    else:
        pdf_source = None

    if file_text["csv_df"] is not None:
        df = file_text["csv_df"]
        csv_preview = df.head(10).to_markdown(index=False)
        csv_context = f"CSV table preview:\n\n{csv_preview}\n\nColumns: {', '.join(df.columns)}"
        csv_source = StringKnowledgeSource(content=csv_context)
    else:
        csv_source = None

    # === Agents ===
    pdf_agent = Agent(
        role="Document Analyst",
        goal="Answer questions using PDF and PPT files only.",
        backstory="You specialize in interpreting content from documents.",
        llm=llm,
        verbose=True,
        knowledge_sources=[pdf_source] if pdf_source else []
    )

    csv_agent = Agent(
        role="Data Analyst",
        goal="Answer questions using CSV data only.",
        backstory="You specialize in interpreting structured tabular data.",
        llm=llm,
        verbose=True,
        knowledge_sources=[csv_source] if csv_source else []
    )

    # === Tasks ===
    doc_task = Task(
        description="Answer user's question using PDF or PPT files.",
        expected_output="Answer derived from document content.",
        agent=pdf_agent,
    )

    csv_task = Task(
        description="Answer user's question using uploaded CSV files.",
        expected_output="Answer derived from structured data.",
        agent=csv_agent,
    )

    # === Crew ===
    crew = Crew(
        agents=[pdf_agent, csv_agent],
        tasks=[doc_task, csv_task],
        process=Process.hierarchical,
        manager_llm=llm,
        verbose=True,
    )

    # === Suggested Questions ===
    def generate_suggestions(context, n=5):
        prompt = f"""Suggest {n} smart questions a user might ask based on this content:\n---\n{context[:3000]}\n---"""
        res = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=300,
        )
        return [line.strip("0123456789. ") for line in res.choices[0].message.content.split("\n") if line.strip()]

    st.divider()
    st.subheader("💬 Ask your question")

    st.markdown("💡 Suggested Questions:")
    for q in generate_suggestions(full_text):
        st.markdown(f"- {q}")

    # Display chat history
    for role, msg in st.session_state.chat_history:
        st.markdown(f"**{role}:** {msg}")

    # Initialize default input state
    if "user_question_input" not in st.session_state:
        st.session_state.user_question_input = ""
        user_query = st.text_input("Type your question", key="user_question_input")

    if st.button("Send") and user_query.strip():
        with st.spinner("Thinking..."):
            response = crew.kickoff(inputs={"question": user_query})

            st.session_state.chat_history.append(("User", user_query))
            st.session_state.chat_history.append(("AI", response))

            st.success("✅ Answer ready!")
            st.markdown(f"**AI:** {response}")
            
            # ✅ Clear the input box by resetting the session state
            st.session_state.user_question_input = ""
